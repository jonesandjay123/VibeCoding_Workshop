#!/usr/bin/env python3
"""Convert P1_slide_content_ja.md (Japanese version) to PowerPoint."""

import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Parse the markdown
with open("/Users/jarvis/Downloads/code/VibeCoding_Workshop/session2/slides/P1_slide_content_ja.md", "r") as f:
    content = f.read()

# Split by "---" to get slide blocks
blocks = content.split("\n---\n")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Colors
BG_DARK = RGBColor(0x1a, 0x1a, 0x2e)
BG_BLOCK = {
    "1": RGBColor(0x2d, 0x1b, 0x69),  # purple
    "2": RGBColor(0x11, 0x44, 0x44),  # teal
    "3": RGBColor(0x44, 0x22, 0x11),  # brown
    "4": RGBColor(0x11, 0x33, 0x55),  # blue
    "5": RGBColor(0x33, 0x44, 0x11),  # green
    "6": RGBColor(0x44, 0x11, 0x33),  # pink
}
ACCENT = RGBColor(0xff, 0xe1, 0x35)  # yellow
WHITE = RGBColor(0xff, 0xff, 0xff)
LIGHT_GRAY = RGBColor(0xcc, 0xcc, 0xcc)
NOTES_GRAY = RGBColor(0xaa, 0xaa, 0xaa)

current_block = "1"

def set_slide_bg(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_text_box(slide, left, top, width, height, text, font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.alignment = alignment
    return tf

for block in blocks:
    block = block.strip()
    if not block:
        continue
    
    # Detect block headers (═══ lines)
    if "═══" in block:
        m = re.search(r'Block (\d)', block)
        if m:
            current_block = m.group(1)
        # Create a section divider slide
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
        set_slide_bg(slide, BG_BLOCK.get(current_block, BG_DARK))
        
        # Extract block title
        lines = [l.strip() for l in block.split('\n') if l.strip() and '═' not in l and l.strip().startswith('#')]
        title_text = lines[0].lstrip('# ') if lines else f"Block {current_block}"
        # Also get subtitle
        subtitle_lines = [l.strip().lstrip('# ') for l in block.split('\n') if l.strip() and '═' not in l and not l.strip().startswith('#') and '⏰' in l]
        subtitle = subtitle_lines[0] if subtitle_lines else ""
        
        add_text_box(slide, 1, 2, 11, 2, title_text, font_size=40, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
        if subtitle:
            add_text_box(slide, 1, 4, 11, 1, subtitle, font_size=24, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
        continue
    
    # Skip non-slide blocks (table of contents, etc.)
    if not block.startswith("## Slide") and not block.startswith("## 付録"):
        # Check if it's the title slide
        if block.startswith("# Session"):
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            set_slide_bg(slide, BG_DARK)
            lines = block.split('\n')
            title = lines[0].lstrip('# ')
            subtitle = lines[1].lstrip('# ') if len(lines) > 1 else ""
            add_text_box(slide, 1, 2, 11, 2, title, font_size=44, color=ACCENT, bold=True, alignment=PP_ALIGN.CENTER)
            add_text_box(slide, 1, 4, 11, 1, subtitle, font_size=28, color=WHITE, alignment=PP_ALIGN.CENTER)
            # Find the quote line
            for l in lines:
                if l.startswith('>'):
                    add_text_box(slide, 1, 5, 11, 1, l.lstrip('> '), font_size=18, color=NOTES_GRAY, alignment=PP_ALIGN.CENTER)
        continue
    
    # Parse slide content
    lines = block.split('\n')
    
    # Get slide title from first ## line
    slide_title = ""
    slide_subtitle = ""
    for line in lines:
        if line.startswith("## Slide") or line.startswith("## 付録"):
            # e.g. "## Slide 1 — オープニング：あなたのサイト、まだ生きてる？"
            parts = line.split("—", 1)
            slide_title = parts[1].strip() if len(parts) > 1 else line.lstrip("## ")
            break
    
    # Get the **標題：** line
    for line in lines:
        if line.startswith("**標題：**"):
            slide_subtitle = line.replace("**標題：**", "").strip()
            break
    
    # Get bullet points (lines starting with -)
    bullets = []
    for line in lines:
        stripped = line.strip()
        if stripped.startswith("- ") and not stripped.startswith("- **講者"):
            bullet_text = stripped.lstrip("- ").replace("**", "").replace("`", "")
            bullets.append(bullet_text)
    
    # Get code blocks
    code_lines = []
    in_code = False
    for line in lines:
        if line.strip().startswith("```"):
            in_code = not in_code
            continue
        if in_code:
            code_lines.append(line)
    
    # Get speaker notes
    notes_text = ""
    in_notes = False
    for line in lines:
        if "> **講者備註" in line or "> 講者備註" in line:
            in_notes = True
            continue
        if in_notes:
            if line.startswith("> "):
                notes_text += line.lstrip("> ") + "\n"
            elif line.strip() == "":
                continue
            else:
                in_notes = False
    
    # Create slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_bg(slide, BG_BLOCK.get(current_block, BG_DARK))
    
    # Title
    display_title = slide_subtitle if slide_subtitle else slide_title
    add_text_box(slide, 0.8, 0.4, 11, 1, display_title, font_size=32, color=ACCENT, bold=True)
    
    # Content area
    y_pos = 1.5
    
    # Bullets
    if bullets:
        bullet_text = "\n".join(f"• {b}" for b in bullets)
        tf = add_text_box(slide, 0.8, y_pos, 7, 4, "", font_size=20, color=WHITE)
        tf.clear()
        for i, b in enumerate(bullets):
            p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
            p.text = f"• {b}"
            p.font.size = Pt(20)
            p.font.color.rgb = WHITE
            p.space_after = Pt(8)
        y_pos += len(bullets) * 0.45
    
    # Code block (if any, show in a box on the right or below)
    if code_lines:
        code_text = "\n".join(code_lines[:15])  # limit
        code_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(max(y_pos, 2.0)), Inches(11), Inches(4)
        )
        ctf = code_box.text_frame
        ctf.word_wrap = True
        cp = ctf.paragraphs[0]
        cp.text = code_text
        cp.font.size = Pt(14)
        cp.font.color.rgb = RGBColor(0x00, 0xff, 0x88)
        cp.font.name = "Menlo"
        # Add dark background to code box
        code_box.fill.solid()
        code_box.fill.fore_color.rgb = RGBColor(0x0d, 0x0d, 0x0d)
    
    # Speaker notes
    if notes_text.strip():
        slide.notes_slide.notes_text_frame.text = notes_text.strip()

# Save
output_path = "/Users/jarvis/Downloads/code/VibeCoding_Workshop/session2/slides/Session2_VibeCoding_JA.pptx"
prs.save(output_path)
print(f"PPT saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
